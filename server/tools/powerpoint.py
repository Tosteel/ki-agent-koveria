from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional

from fastapi import HTTPException

from server.services.llm_ionos import IonosLLM


def _wrap_text(text: str, max_chars: int) -> List[str]:
    words = (text or "").split()
    if not words:
        return []

    lines: List[str] = []
    cur = words[0]
    for w in words[1:]:
        if len(cur) + 1 + len(w) <= max_chars:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines


def _split_paragraph_into_lines(paragraph: str, max_chars_per_line: int) -> List[str]:
    txt = (paragraph or "").strip()
    if not txt:
        return []

    # Keep lightweight section markers readable on their own line.
    if txt.startswith("- "):
        body = txt[2:].strip()
        lines = _wrap_text(body, max_chars_per_line - 2)
        if not lines:
            return ["-"]
        return [f"- {lines[0]}"] + [f"  {ln}" for ln in lines[1:]]

    return _wrap_text(txt, max_chars_per_line)


def _group_lines_to_boxes(lines: List[str], max_chars_per_box: int, max_lines_per_box: int) -> List[str]:
    boxes: List[str] = []
    cur: List[str] = []
    cur_chars = 0
    for ln in lines:
        add = len(ln) + (1 if cur else 0)
        too_many_chars = cur and (cur_chars + add > max_chars_per_box)
        too_many_lines = cur and (len(cur) >= max_lines_per_box)
        if too_many_chars or too_many_lines:
            boxes.append("\n".join(cur).strip())
            cur = []
            cur_chars = 0
        cur.append(ln)
        cur_chars += len(ln) + 1
    if cur:
        boxes.append("\n".join(cur).strip())
    return [b for b in boxes if b]


def layout_text_for_slides(
    title: str,
    text: str,
    *,
    max_chars_per_line: int = 88,
    max_chars_per_slide: int = 900,
    max_lines_per_slide: int = 14,
    max_chars_per_box: int = 330,
    max_lines_per_box: int = 5,
) -> List[Dict[str, object]]:
    """
    Split plain text into slide-friendly blocks.
    Returns list of slides with shape:
      {"title": str, "boxes": List[str]}
    """
    raw = (text or "").strip()
    if not raw:
        return [{"title": title or "Ergebnis", "boxes": ["Kein Inhalt."]}]

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", raw) if p.strip()]
    para_lines: List[str] = []
    for p in paragraphs:
        para_lines.extend(_split_paragraph_into_lines(p, max_chars_per_line))
        para_lines.append("")  # spacer between paragraphs
    if para_lines and para_lines[-1] == "":
        para_lines.pop()

    all_boxes = _group_lines_to_boxes(para_lines, max_chars_per_box=max_chars_per_box, max_lines_per_box=max_lines_per_box)
    if not all_boxes:
        all_boxes = ["Kein Inhalt."]

    slides: List[Dict[str, object]] = []
    cur_boxes: List[str] = []
    cur_chars = 0
    cur_lines = 0

    for box in all_boxes:
        box_lines = [ln for ln in box.splitlines()]
        box_chars = len(box)
        would_exceed = (
            cur_boxes
            and (cur_chars + box_chars > max_chars_per_slide or cur_lines + len(box_lines) > max_lines_per_slide)
        )
        if would_exceed:
            idx = len(slides) + 1
            ttl = title if idx == 1 else f"{title} (Teil {idx})"
            slides.append({"title": ttl or "Ergebnis", "boxes": cur_boxes})
            cur_boxes = []
            cur_chars = 0
            cur_lines = 0

        cur_boxes.append(box)
        cur_chars += box_chars
        cur_lines += len(box_lines)

    if cur_boxes:
        idx = len(slides) + 1
        ttl = title if idx == 1 else f"{title} (Teil {idx})"
        slides.append({"title": ttl or "Ergebnis", "boxes": cur_boxes})

    return slides


def _parse_json_strictish(text: str) -> Dict[str, Any]:
    if not text:
        return {}
    t = text.strip()
    if t.startswith("```"):
        t = t.strip("`").strip()
        if "\n" in t:
            first, rest = t.split("\n", 1)
            if first.strip().lower() in ("json", "javascript"):
                t = rest.strip()
    try:
        obj = json.loads(t)
        return obj if isinstance(obj, dict) else {}
    except Exception:
        pass

    start = t.find("{")
    end = t.rfind("}")
    if start >= 0 and end > start:
        snippet = t[start : end + 1]
        try:
            obj = json.loads(snippet)
            return obj if isinstance(obj, dict) else {}
        except Exception:
            return {}
    return {}


def _normalize_llm_slides(
    data: Dict[str, Any],
    fallback_title: str,
    *,
    max_slides: int,
    max_boxes_per_slide: int,
) -> List[Dict[str, object]]:
    slides_raw = data.get("slides")
    if not isinstance(slides_raw, list):
        alt = data.get("presentation")
        if isinstance(alt, dict):
            slides_raw = alt.get("slides")
    if not isinstance(slides_raw, list):
        alt_sections = data.get("sections")
        if isinstance(alt_sections, list):
            slides_raw = alt_sections
    if not isinstance(slides_raw, list):
        return []

    slides: List[Dict[str, object]] = []
    for s in slides_raw:
        if not isinstance(s, dict):
            continue
        title = str(s.get("title") or fallback_title or "Ergebnis").strip()[:140]
        boxes_raw = s.get("boxes")
        boxes: List[str] = []
        if isinstance(boxes_raw, list):
            boxes = [str(b).strip() for b in boxes_raw if str(b).strip()]
        elif isinstance(boxes_raw, str) and boxes_raw.strip():
            boxes = [boxes_raw.strip()]
        boxes = boxes[:max_boxes_per_slide]
        if not boxes:
            continue
        slides.append({"title": title, "boxes": boxes})
        if len(slides) >= max_slides:
            break
    return slides


def _layout_text_for_slides_llm_with_mode(
    title: str,
    text: str,
    *,
    goal: str = "",
    instruction: str = "",
    max_slides: int = 12,
    max_boxes_per_slide: int = 3,
    llm: Optional[IonosLLM] = None,
) -> tuple[List[Dict[str, object]], str, str]:
    raw = (text or "").strip()
    if not raw:
        return ([{"title": title or "Ergebnis", "boxes": ["Kein Inhalt."]}], "heuristic", "empty_input")

    client = llm or IonosLLM()
    if not client.enabled():
        return (layout_text_for_slides(title=title, text=text), "heuristic", "llm_not_enabled")

    system = (
        "Du bist ein Präsentations-Layout-Assistent.\n"
        "Erzeuge NUR JSON mit Struktur: {\"slides\":[{\"title\":\"...\",\"boxes\":[\"...\",\"...\"]}]}\n"
        "Regeln:\n"
        "- Jede Slide hat 1 Titel und 1-3 Textboxen.\n"
        "- Jede Textbox enthält prägnante, gut lesbare Textblöcke.\n"
        "- Keine Meta-Sätze, keine Hinweise auf den Erstellungsprozess.\n"
        "- Nur Inhalte aus dem Input verwenden.\n"
    )
    if instruction.strip():
        system += f"Zusatzanweisung: {instruction.strip()}\n"

    user = (
        f"Ziel: {goal.strip() or '(nicht angegeben)'}\n"
        f"Präsentationstitel: {title or 'Ergebnis'}\n"
        f"Maximale Anzahl Slides: {int(max_slides)}\n"
        f"Maximale Textboxen pro Slide: {int(max_boxes_per_slide)}\n\n"
        f"Inhalt:\n{raw}"
    )

    schema = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "title": {"type": "string"},
                        "boxes": {
                            "type": "array",
                            "items": {"type": "string"},
                            "minItems": 1,
                            "maxItems": max(1, int(max_boxes_per_slide)),
                        },
                    },
                    "required": ["title", "boxes"],
                },
                "minItems": 1,
                "maxItems": max(1, int(max_slides)),
            }
        },
        "required": ["slides"],
    }
    response_format = {
        "type": "json_schema",
        "json_schema": {
            "name": "ppt_layout",
            "schema": schema,
            "strict": True,
        },
    }

    completion = client.chat_completions(
        messages=[{"role": "system", "content": system}, {"role": "user", "content": user}],
        max_tokens=1800,
        temperature=0.0,
        top_p=0.1,
        response_format=response_format,
    )
    parsed = _parse_json_strictish(client.extract_text(completion))
    if not parsed:
        # Fallback auf json_object, falls json_schema vom Provider nicht strikt eingehalten wird.
        completion2 = client.chat_completions(
            messages=[{"role": "system", "content": system}, {"role": "user", "content": user}],
            max_tokens=1800,
            temperature=0.0,
            top_p=0.1,
            response_format={"type": "json_object"},
        )
        parsed = _parse_json_strictish(client.extract_text(completion2))

    slides = _normalize_llm_slides(
        parsed,
        title or "Ergebnis",
        max_slides=max_slides,
        max_boxes_per_slide=max_boxes_per_slide,
    )
    if slides:
        return (slides, "llm", "")
    return (layout_text_for_slides(title=title, text=text), "heuristic", "llm_output_invalid_or_empty")


def layout_text_for_slides_llm(
    title: str,
    text: str,
    *,
    goal: str = "",
    instruction: str = "",
    max_slides: int = 12,
    max_boxes_per_slide: int = 3,
    llm: Optional[IonosLLM] = None,
) -> List[Dict[str, object]]:
    slides, _mode, _reason = _layout_text_for_slides_llm_with_mode(
        title=title,
        text=text,
        goal=goal,
        instruction=instruction,
        max_slides=max_slides,
        max_boxes_per_slide=max_boxes_per_slide,
        llm=llm,
    )
    return slides


def export_text_pptx(
    output_file: Path,
    title: str,
    text: str,
    *,
    use_llm_layout: bool = False,
    allow_heuristic_fallback: bool = True,
    goal: str = "",
    instruction: str = "",
    max_slides: int = 12,
    max_boxes_per_slide: int = 3,
    llm: Optional[IonosLLM] = None,
) -> Dict[str, Any]:
    output_file.parent.mkdir(parents=True, exist_ok=True)
    if output_file.suffix.lower() != ".pptx":
        raise HTTPException(status_code=400, detail="output_path must end with .pptx")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        raise HTTPException(status_code=500, detail="python-pptx not installed. Please install package 'python-pptx'.")

    prs = Presentation()
    if use_llm_layout:
        slides, layout_mode, fallback_reason = _layout_text_for_slides_llm_with_mode(
            title=title,
            text=text,
            goal=goal,
            instruction=instruction,
            max_slides=max_slides,
            max_boxes_per_slide=max_boxes_per_slide,
            llm=llm,
        )
        if layout_mode != "llm" and not allow_heuristic_fallback:
            raise HTTPException(
                status_code=502,
                detail=f"LLM layout failed ({fallback_reason}). Set allow_heuristic_fallback=true to continue with heuristic layout.",
            )
    else:
        slides = layout_text_for_slides(title=title, text=text)
        layout_mode = "heuristic"

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), slide_w - Inches(1.0), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.clear()
        p = tf_title.paragraphs[0]
        p.text = str(slide_data.get("title") or "Ergebnis")
        p.font.size = Pt(28)
        p.font.bold = True

        boxes = slide_data.get("boxes") or []
        if not isinstance(boxes, list):
            boxes = [str(boxes)]

        top = Inches(1.1)
        bottom_margin = Inches(0.4)
        usable_h = max(Inches(1.0), slide_h - top - bottom_margin)
        count = max(1, len(boxes))
        gap = Inches(0.15)
        each_h = (usable_h - gap * (count - 1)) / count

        for i, box_text in enumerate(boxes):
            y = top + i * (each_h + gap)
            tb = slide.shapes.add_textbox(Inches(0.7), y, slide_w - Inches(1.4), each_h)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(box_text)
            p.font.size = Pt(16)

    prs.save(str(output_file))
    return {
        "bytes_written": output_file.stat().st_size,
        "layout_mode": layout_mode,
    }
